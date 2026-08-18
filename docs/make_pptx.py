# -*- coding: utf-8 -*-
"""生成《赛菲莉娅背包整理插件 工作原理》PPT（16:9）。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, 'figures')
OUT = os.path.join(BASE, '赛菲莉娅背包整理插件-工作原理.pptx')

DARK = RGBColor(0x1F, 0x3B, 0x57)
BLUE = RGBColor(0x2F, 0x6F, 0xAD)
GRAY = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x3D, 0x8B, 0x4F)
ORANGE = RGBColor(0xC4, 0x7A, 0x1F)
RED = RGBColor(0xB4, 0x43, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_font(run, size, bold=False, color=DARK, name='微软雅黑'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def add_text(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.12):
    """lines: list of (text, size, bold, color) 或 (text, size, bold, color, space_before)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in lines:
        text, size, bold, color = item[0], item[1], item[2], item[3]
        space = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if space:
            p.space_before = Pt(space)
        r = p.add_run()
        r.text = text
        set_font(r, size, bold, color)
    return tb


def title_bar(slide, text, subtitle=None):
    # 左侧色块
    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.42), Inches(0.09), Inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    lines = [(text, 26, True, DARK)]
    add_text(slide, 0.85, 0.36, 11.9, 0.8, lines)
    if subtitle:
        add_text(slide, 0.85, 1.02, 11.9, 0.4, [(subtitle, 12.5, False, GRAY)])
    # 分隔线
    ln = slide.shapes.add_shape(1, Inches(0.55), Inches(1.42), Inches(12.23), Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor(0xD8, 0xE2, 0xEC)
    ln.line.fill.background()


def bullet_box(slide, left, top, width, height, items, size=15, gap=6, color=DARK):
    """items: list of str 或 (str, sub)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            main, sub = it
        else:
            main, sub = it, None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = '· ' + main
        set_font(r, size, False, color)
        if sub:
            r2 = p.add_run()
            r2.text = '  ' + sub
            set_font(r2, size - 2.5, False, LGRAY)
    return tb


def place_pic(slide, fname, left_in, top_in, max_w_in, max_h_in):
    path = os.path.join(FIG, fname)
    im = Image.open(path)
    w, h = im.size
    scale = min(max_w_in * 96 / w, max_h_in * 96 / h)
    w_in = w * scale / 96
    h_in = h * scale / 96
    left = left_in + (max_w_in - w_in) / 2 if left_in >= 0 else (13.333 - w_in) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top_in), Inches(w_in), Inches(h_in))
    return w_in, h_in


def new_slide():
    return prs.slides.add_slide(BLANK)


# ================================================================ P1 封面
s = new_slide()
# 顶部蓝色色带
band = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
add_text(s, 1.2, 2.0, 11.0, 1.0, [('赛菲莉娅 背包整理插件', 44, True, DARK)])
add_text(s, 1.2, 3.1, 11.0, 0.8, [('工作原理：一键 F8 智能整理背包', 26, True, BLUE)])
add_text(s, 1.2, 4.15, 11.0, 0.6, [('识别 → 评分 → 搜索 → 应用', 17, False, GRAY)])
add_text(s, 1.2, 5.9, 11.0, 0.5, [('BepInEx 插件 · v2.3.9 · Enhanced 增强整理模式', 13, False, LGRAY)])
add_text(s, 1.2, 6.35, 11.0, 0.5, [('2026-08', 13, False, LGRAY)])
# 底部装饰
band2 = s.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2))
band2.fill.solid()
band2.fill.fore_color.rgb = RGBColor(0xD8, 0xE2, 0xEC)
band2.line.fill.background()

# ================================================================ P2 解决什么问题
s = new_slide()
title_bar(s, '它解决什么问题', '手动摆背包，几乎不可能同时满足所有加成')
bullet_box(s, 0.7, 1.75, 6.1, 4.6, [
    ('石板改格子等级：', '+N 等级 / ×M 倍率 / 解锁豁免格'),
    ('护符带位置条件：', '顶行、两侧、内侧、旁边必须空着、八邻域全满……'),
    ('组合机制要"摆在一起"：', '行星望远镜+行星、和谐之晶周围8格、奉献徽章同行、指北针上下配对'),
    ('游戏内置排列只跑几次迭代，', '经常停在明显不是最优的结果'),
])
box = s.shapes.add_shape(1, Inches(7.1), Inches(1.9), Inches(5.5), Inches(4.3))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)
box.line.color.rgb = BLUE
box.line.width = Pt(1.2)
add_text(s, 7.45, 2.15, 4.9, 3.8, [
    ('插件的做法', 18, True, BLUE),
    ('按 F8，把背包重排一遍，', 15, False, DARK, 10),
    ('让各种加成尽量同时吃到。', 15, False, DARK, 2),
    ('排完提示"整理完毕"。', 15, False, DARK, 2),
    ('不修改游戏文件，', 12.5, False, GRAY, 12),
    ('单机 / 主机 / 联机客户端都可用。', 12.5, False, GRAY, 2),
])

# ================================================================ P3 总体思路
s = new_slide()
title_bar(s, '总体思路：先离线算，再一次性应用', '在背包副本上算最优，游戏里只改一次')
bullet_box(s, 0.7, 1.75, 5.4, 3.4, [
    ('1. 识别', '物品分类 + 特殊机制识别'),
    ('2. 评分', '给任意摆法打分（越高加成越大）'),
    ('3. 搜索', '智能初始 + 多轮模拟退火找最高分'),
], size=15)
bullet_box(s, 0.7, 4.35, 5.4, 2.6, [
    ('副本上可跑几千次迭代', '游戏内置算法迭代次数少（默认仅 4 次）'),
    ('游戏内只改一次，干净利落', ''),
    ('联机客户端也能用', '算完用网络命令执行'),
], size=14)
place_pic(s, 'fig_overview.png', 6.1, 1.8, 6.7, 4.9)

# ================================================================ P4 识别层
s = new_slide()
title_bar(s, '第一步：认识背包里的每一件东西', '识别靠类型/接口判断，配置 key 兜底')
bullet_box(s, 0.7, 1.7, 5.9, 2.4, [
    ('石板：', '占多格 + 效果网格（加等级/乘倍率/解锁豁免格），可旋转'),
    ('护符：', '占一格 + 等级 + 位置条件/联动机制'),
    ('负面藏品（负担）：', '心之重担等，塞进最差的格子减少损失'),
    ('杂物：', '凑数物品，填空'),
], size=13.5, gap=5)
add_text(s, 0.7, 4.05, 5.9, 0.4, [('特殊机制识别', 15, True, BLUE)])
bullet_box(s, 0.7, 4.5, 5.9, 2.6, [
    ('位置条件护符（9 种枚举） / 豁免格（解锁石板）', ''),
    ('行星望远镜 + 行星藏品（PLANET 标签，可排除）', ''),
    ('和谐之晶 / 奉献徽章 / 同伴（ICompanionCharm 接口）', ''),
    ('指北针 / 行锁定物品（凯尔萨德尼钥匙）', ''),
    ('神秘 ×2 地块 / 武器相关护符（需武器匹配）', ''),
], size=13, gap=5)
add_text(s, 6.9, 1.7, 5.9, 0.4, [('一张布局能同时吃到的机制', 15, True, BLUE)])
place_pic(s, 'fig_backpack.png', 6.9, 2.2, 5.9, 4.9)

# ================================================================ P5 评分模型
s = new_slide()
title_bar(s, '第二步：给一种摆法打分', '离线模拟游戏全部加成公式，作为搜索的目标函数')
bullet_box(s, 0.7, 1.75, 6.1, 3.3, [
    ('有效等级 =', '石板贡献 + 附魔，神秘地块再 ×2'),
    ('是否启用：', '等级≥0 && 位置满足（或豁免格）&& 武器匹配'),
    ('核心分：', '有效等级 × 10000 × 优先级权重'),
    ('启用 +1000 / 禁用 -750，负格 -250/级', ''),
    ('优先级：', '传说/羁绊 P1 → 普通 P4，权重 1.5 → 1.0'),
], size=14, gap=6)
add_text(s, 0.7, 5.05, 6.1, 0.4, [('机制联动分（默认权重）', 14, True, BLUE)])
bullet_box(s, 0.7, 5.45, 6.1, 1.8, [
    ('行星聚簇 +40000/颗 · 和谐之晶 +2000×等级和', ''),
    ('奉献徽章 +3000/同伴 · 指北针配对 +12000', ''),
    ('负担惩罚 -20000/级 · 行锁定跨行 -100000', ''),
], size=13, gap=4)
place_pic(s, 'fig_weights.png', 6.9, 1.8, 6.0, 5.0)

# ================================================================ P6 等级分配
s = new_slide()
title_bar(s, '石板和物品怎么配合：等级分配', '等级是格子的属性——石板决定哪几格亮，护符决定站哪格')
bullet_box(s, 0.7, 1.8, 5.9, 4.9, [
    ('格子等级 =', '(裸等级 + 石板贡献) × 神秘倍率 + 附魔'),
    ('石板 = 聚光灯', '位置/旋转决定哪几格 +N、哪格禁用/豁免'),
    ('护符 = 演员', '位置条件限区域，其余按优先级挑亮格'),
    ('"先石板后物品"是交替推进', '每摆一步重算全部格子等级，下一步基于最新等级'),
    ('负效果石板最先摆', '负等级不漏到将来放护符的格子上'),
    ('退火里旋转石板 + 移动物品', '灯光与站位联合优化，直到评分收敛'),
], size=14, gap=8)
place_pic(s, 'fig_levels.png', 6.7, 1.9, 6.3, 4.9)

# ================================================================ P7 简单案例
s = new_slide()
title_bar(s, '简单案例：同一块石板，两种摆法', '只挪"灯光"和"站位"，总评分 65000 → 95000（+46%）')
place_pic(s, 'fig_case.png', 0.6, 1.9, 12.1, 4.9)
add_text(s, 0.7, 6.85, 12.0, 0.5, [
    ('整理前 P1 站 1 级格，拿 16000 分；整理后 P1 站到石板加成的 3 级格，拿 46000 分——'
     '评分公式：有效等级 × 10000 × 优先级权重 + 启用 1000（P1 传说权重 1.5）', 11.5, False, GRAY),
], align=PP_ALIGN.CENTER)

# ================================================================ P6 智能初始布局
s = new_slide()
title_bar(s, '第三步：搜索（一）智能初始布局', '不靠运气，按机制优先级逐个落位')
bullet_box(s, 0.7, 1.75, 5.6, 5.2, [
    ('石板贪心摆位', '负效果优先，逐格逐旋转打分'),
    ('受限护符先放', '满足条件格 → 豁免格 → 任意格'),
    ('行锁定物品', '保持用户所在行，行内选最佳列'),
    ('望远镜 / 和谐之晶 / 奉献徽章', '先落位：8 邻域加权、同行加权引导后续'),
    ('行星聚簇 → 罗盘配对', '望远镜旁 / "上方是伤害类"的格子'),
    ('其余护符按优先级，负担塞负格', ''),
], size=14, gap=7)
place_pic(s, 'fig_smartstart.png', 6.4, 1.75, 6.5, 5.1)

# ================================================================ P7 模拟退火
s = new_slide()
title_bar(s, '第三步：搜索（二）模拟退火', '在起点附近反复试探，温度从高到低收敛')
bullet_box(s, 0.7, 1.75, 12.0, 1.7, [
    ('每次迭代做一次小改动并重新打分：约 38% 是定向移动（受限护符跳满足格 / 行星挪到望远镜旁 / 指北针挪到伤害类下方 / 负担丢进负格），'
     '其余是随机探索（移动 / 交换 / 旋转石板）'),
    ('温度高时接受"变差"的解，翻过局部最优的小山坡；温度线性降到 1 后只接受变好，保证收敛；每轮随机重启 3 次'),
], size=14, gap=8)
place_pic(s, 'fig_anneal.png', 0.7, 3.6, 12.0, 3.3)

# ================================================================ P8 多轮搜索
s = new_slide()
title_bar(s, '第三步：搜索（三）多轮独立搜索', 'SearchRounds = 4：每按一次 F8 跑 4 轮，取全局最优')
bullet_box(s, 0.7, 1.85, 6.0, 4.6, [
    ('每轮不同随机种子', '等效自动把 F8 重复按 4 次'),
    ('多起点', '智能初始 / 原始布局 / 随机打乱'),
    ('取全局最高分', '一次按键直达最佳'),
    ('全离线评估，毫秒级', '34 格满包约 200~300ms'),
    ('分数不升则不应用', '结果不比现状好就保持原样'),
], size=15, gap=10)
place_pic(s, 'fig_mutation.png', 6.9, 1.85, 5.9, 4.9)

# ================================================================ P9 主机 vs 客户端
s = new_slide()
title_bar(s, '第四步：把结果写回游戏', '单机/主机与联机客户端两种执行方式')
# 左：主机
box = s.shapes.add_shape(1, Inches(0.7), Inches(1.8), Inches(5.9), Inches(2.6))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xED, 0xF3, 0xFA)
box.line.color.rgb = BLUE
box.line.width = Pt(1)
add_text(s, 1.0, 2.0, 5.3, 2.2, [
    ('单机 / 主机', 17, True, BLUE),
    ('直接把算好的整包布局写回背包', 14, False, DARK, 8),
    ('再用游戏接口读一次真实评分核对', 14, False, DARK, 3),
    ('日志对比"离线评分 / 游戏评分"', 12.5, False, GRAY, 6),
])
# 右：客户端
box = s.shapes.add_shape(1, Inches(6.9), Inches(1.8), Inches(5.9), Inches(2.6))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF3, 0xEF, 0xE7)
box.line.color.rgb = ORANGE
box.line.width = Pt(1)
add_text(s, 7.2, 2.0, 5.3, 2.2, [
    ('联机客户端', 17, True, ORANGE),
    ('没有服务器权限，改用网络命令', 14, False, DARK, 8),
    ('"当前布局 → 目标布局"翻译成交换 / 旋转序列', 14, False, DARK, 3),
    ('逐个执行，不清空背包，更保守', 12.5, False, GRAY, 6),
])
place_pic(s, 'fig_client.png', 0.7, 4.6, 12.0, 2.6)

# ================================================================ P10 安全设计
s = new_slide()
title_bar(s, '安全设计：宁可不动，不可丢物品', '整理是重排整个背包，错误时机动手可能丢东西')
cards = [
    ('进图等 3 秒', 'SessionStableDelay：联机房间刚进时背包初始化未完成，此时按 F8 提示"请稍候"，不动手', BLUE),
    ('快照一致性校验', '整理前抓背包快照，只统计正常背包格（排除药水带等特殊区域），对不上就取消本次整理', GREEN),
    ('一次应用，失败不重试', '算完只应用一次，宁可留着不整理也不冒险', ORANGE),
]
x = 0.7
for title, body, color in cards:
    box = s.shapes.add_shape(1, Inches(x), Inches(1.95), Inches(3.9), Inches(3.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFB)
    box.line.color.rgb = color
    box.line.width = Pt(1.4)
    add_text(s, x + 0.3, 2.2, 3.3, 2.6, [
        (title, 17, True, color),
        (body, 13, False, DARK, 10),
    ])
    x += 4.13
add_text(s, 0.7, 5.3, 12.0, 1.3, [
    ('这些防御是有代价换来的：早期版本用"清空再重写"应用布局，触发过游戏起始物品补货机制，出现过物品丢失 / 复制。'
     '后来改成防御式设计，整理前多重校验。', 14, False, GRAY),
])

# ================================================================ P11 实测与日志
s = new_slide()
title_bar(s, '实测效果与日志诊断', '每次整理都输出完整落地情况，方便核对')
bullet_box(s, 0.7, 1.8, 5.9, 3.0, [
    ('一次 F8 直达最佳', '多轮搜索合并，单键即达最优'),
    ('和谐之晶实测', '3 晶聚簇 + 14 级护符，评分 4000 → 355004，141ms'),
    ('40 件满包约 180ms', '毫秒级，游戏内无卡顿感'),
], size=14, gap=8)
add_text(s, 0.7, 4.6, 5.9, 0.4, [('日志输出（示意）', 14, True, BLUE)])
logbox = s.shapes.add_shape(1, Inches(0.7), Inches(5.0), Inches(5.9), Inches(1.9))
logbox.fill.solid()
logbox.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
logbox.line.fill.background()
add_text(s, 1.0, 5.15, 5.3, 1.6, [
    ('智能初始布局：4000 -> 210000', 11, False, RGBColor(0xD8, 0xD8, 0xD8)),
    ('增强整理完成（141ms）：4000 -> 355004', 11, False, RGBColor(0xD8, 0xD8, 0xD8)),
    ('奉献徽章@0,5：同行同伴 2 个', 11, False, RGBColor(0xD8, 0xD8, 0xD8)),
])
add_text(s, 6.9, 1.9, 6.0, 0.4, [('每次整理输出的诊断信息', 15, True, BLUE)])
bullet_box(s, 6.9, 2.4, 6.0, 4.5, [
    ('识别统计：护符 / 望远镜 / 罗盘 / 附魔 / 奉献徽章 / 同伴数量', ''),
    ('布局网格图：每格物品类型 + 等级', ''),
    ('机制落地位置：', ''),
    ('望远镜@x,y 相邻行星 N 颗', ''),
    ('和谐之晶@x,y 周围 8 格等级和', ''),
    ('奉献徽章@x,y 同行同伴 N 个', ''),
    ('罗盘@x,y 上方类型 / 负担@x,y 格等级', ''),
    ('离线评分 / 游戏评分前后对比', ''),
], size=13, gap=6)

# ================================================================ P12 配置速查
s = new_slide()
title_bar(s, '配置速查', '全部参数在 com.sephiria.backpack-organizer.cfg，改完重启生效')
rows = [
    ('General / Hotkey', 'F8', '触发整理的快捷键'),
    ('General / SessionStableDelay', '3 秒', '进图后等待背包初始化（防丢物品）'),
    ('Enhanced / SearchRounds', '4', '每按一次 F8 跑的独立搜索轮数'),
    ('Enhanced / Iterations · Temperature', '3000 · 800', '退火迭代次数与初始温度'),
    ('Priority / Weight1~4', '1.5 / 1.25 / 1.1 / 1.0', '各优先级护符的等级分权重'),
    ('Synergy / PlanetBonus', '40000', '行星聚簇：望远镜周围每颗行星'),
    ('Synergy / HarmonyLevelBonus', '2000', '和谐之晶：周围 8 格每级等级'),
    ('Synergy / DedicationCompanionBonus', '3000', '奉献徽章：每个同行同伴'),
    ('Synergy / CompassBonus', '12000', '指北针配对奖励'),
    ('Burden / NegativeCellPenalty', '20000', '负担未待负格时的扣分'),
]
top = 1.75
row_h = 0.44
for i, (k, v, d) in enumerate(rows):
    y = top + i * row_h
    if i % 2 == 0:
        box = s.shapes.add_shape(1, Inches(0.7), Inches(y), Inches(12.0), Inches(row_h))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFA)
        box.line.fill.background()
    add_text(s, 0.95, y + 0.08, 4.4, 0.34, [(k, 12.5, True, DARK)])
    add_text(s, 5.5, y + 0.08, 2.6, 0.34, [(v, 12.5, True, BLUE)])
    add_text(s, 8.3, y + 0.08, 4.3, 0.34, [(d, 12, False, GRAY)])

# ================================================================ P13 结束
s = new_slide()
band = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
add_text(s, 1.2, 2.2, 11.0, 0.9, [('一句话总结', 20, True, BLUE)])
add_text(s, 1.2, 3.1, 11.0, 0.9, [('识别每件物品 → 给摆法打分 → 搜索最高分布局 → 一次性应用', 26, True, DARK)])
add_text(s, 1.2, 4.4, 11.0, 0.6, [('按 F8，剩下的交给插件。', 17, False, GRAY)])
add_text(s, 1.2, 5.6, 11.0, 0.5, [('BepInEx 插件 · v2.3.9 · 单机 / 主机 / 联机客户端', 13, False, LGRAY)])
band2 = s.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2))
band2.fill.solid()
band2.fill.fore_color.rgb = RGBColor(0xD8, 0xE2, 0xEC)
band2.line.fill.background()

prs.save(OUT)
print('pptx saved ->', OUT)
