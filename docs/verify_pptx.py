# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

path = sys.argv[1]
p = Presentation(path)
print('pages:', len(p.slides))
for i, s in enumerate(p.slides, 1):
    pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
    texts = [sh.text_frame.text.strip().replace('\n', ' / ')
             for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    t0 = texts[0][:34] if texts else '(no text)'
    t1 = texts[1][:34] if len(texts) > 1 else ''
    print('P%d: pics=%d | %s | %s' % (i, pics, t0, t1))
